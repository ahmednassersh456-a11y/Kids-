from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

BLANK = 6  # blank slide layout index

def rgb(r, g, b):
    return RGBColor(r, g, b)

def add_slide(prs):
    layout = prs.slide_layouts[BLANK]
    return prs.slides.add_slide(layout)

def set_bg(slide, r, g, b):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb(r, g, b)

def add_textbox(slide, text, left, top, width, height, font_size, bold=False,
                color=(0,0,0), align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(*color)
    return txBox

def add_rect(slide, left, top, width, height, fill_color, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(*fill_color)
    if line_color:
        shape.line.color.rgb = rgb(*line_color)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        5,  # rounded rectangle
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(*fill_color)
    if line_color:
        shape.line.color.rgb = rgb(*line_color)
    else:
        shape.line.fill.background()
    return shape

# ==================== SLIDE 1: TITLE ====================
s1 = add_slide(prs)
set_bg(s1, 255, 245, 200)

add_rounded_rect(s1, 0.4, 0.2, 9.2, 1.3, (255, 180, 50))
add_textbox(s1, "🧠 Brain Boosters for Kids!", 0.4, 0.25, 9.2, 1.0,
            44, bold=True, color=(255, 255, 255), align=PP_ALIGN.CENTER)
add_textbox(s1, "Ages 3 – 5", 3.8, 1.55, 2.4, 0.4,
            18, bold=True, color=(180, 80, 0), align=PP_ALIGN.CENTER)

add_rounded_rect(s1, 0.6, 2.1, 8.8, 4.8, (255, 255, 255))
add_textbox(s1, "✨ What's Inside ✨", 2.5, 2.2, 5.0, 0.5,
            22, bold=True, color=(80, 30, 0), align=PP_ALIGN.CENTER)

col1 = "🍎  Count the Fruits\n🔷  Match the Shapes\n🔤  Find the Letter B\n🦁  Odd One Out"
col2 = "✏️  Trace the Letters\n📏  Tall or Short?\n🔢  Connect the Dots\n🔡  Fill the Letter\n🌟  Star Certificate"
add_textbox(s1, col1, 0.9, 2.85, 4.0, 2.8, 14, color=(60, 60, 120))
add_textbox(s1, col2, 5.1, 2.85, 4.0, 2.8, 14, color=(60, 60, 120))
add_textbox(s1, "Fun Pre-K Learning Activities  🌈  Ages 3–5", 1.5, 6.9, 7.0, 0.4,
            11, color=(150, 100, 0), align=PP_ALIGN.CENTER)

# ==================== SLIDE 2: COUNT THE FRUITS ====================
s2 = add_slide(prs)
set_bg(s2, 220, 255, 220)

add_rounded_rect(s2, 0.4, 0.2, 9.2, 0.9, (100, 200, 100))
add_textbox(s2, "🍎 Count the Fruits!", 0.4, 0.22, 9.2, 0.85,
            36, bold=True, color=(255, 255, 255), align=PP_ALIGN.CENTER)
add_textbox(s2, "How many do you see? Write the number in the box!",
            0.8, 1.25, 8.4, 0.5, 16, color=(60, 120, 60), align=PP_ALIGN.CENTER)

fruits = [
    ("🍎🍎🍎🍎", "Apples", "4"),
    ("🍌🍌🍌", "Bananas", "3"),
    ("🍓🍓🍓🍓🍓", "Berries", "5"),
    ("🍇🍇", "Grapes", "2"),
]
for i, (emojis, name, _) in enumerate(fruits):
    x = 0.5 + i * 2.35
    add_rounded_rect(s2, x, 1.9, 2.1, 3.0, (255, 255, 255))
    add_textbox(s2, emojis, x + 0.1, 2.0, 1.9, 1.4, 22, align=PP_ALIGN.CENTER)
    add_textbox(s2, name, x + 0.1, 3.45, 1.9, 0.4, 13, color=(80, 80, 80), align=PP_ALIGN.CENTER)
    add_rounded_rect(s2, x + 0.45, 3.9, 1.2, 0.65, (240, 240, 240))
    add_textbox(s2, "____", x + 0.45, 3.95, 1.2, 0.55, 18, color=(100, 100, 100), align=PP_ALIGN.CENTER)

add_textbox(s2, "🌟 Fun Learning for Little Ones  🌈  Pre-K",
            1.5, 6.9, 7.0, 0.4, 11, color=(100, 160, 100), align=PP_ALIGN.CENTER)

# ==================== SLIDE 3: MATCH THE SHAPES ====================
s3 = add_slide(prs)
set_bg(s3, 220, 235, 255)

add_rounded_rect(s3, 0.4, 0.2, 9.2, 0.9, (80, 100, 220))
add_textbox(s3, "🔷 Match the Shapes!", 0.4, 0.22, 9.2, 0.85,
            36, bold=True, color=(255, 255, 255), align=PP_ALIGN.CENTER)
add_textbox(s3, "Draw a line from each shape to its correct name!",
            0.8, 1.25, 8.4, 0.5, 16, color=(60, 60, 160), align=PP_ALIGN.CENTER)

left_shapes  = ["🔵", "🟥", "🔺", "▬", "⭐"]
right_labels = ["Triangle", "Star", "Rectangle", "Circle", "Square"]
for i in range(5):
    y = 1.9 + i * 0.9
    add_rounded_rect(s3, 0.8, y, 1.6, 0.72, (200, 210, 255))
    add_textbox(s3, left_shapes[i], 0.85, y + 0.08, 1.5, 0.55, 26, align=PP_ALIGN.CENTER)
    add_rounded_rect(s3, 7.6, y, 1.7, 0.72, (200, 210, 255))
    add_textbox(s3, right_labels[i], 7.65, y + 0.16, 1.6, 0.4,
                14, bold=True, color=(60, 60, 160), align=PP_ALIGN.CENTER)

add_textbox(s3, "🌟 Fun Learning for Little Ones  🌈  Pre-K",
            1.5, 6.9, 7.0, 0.4, 11, color=(80, 80, 180), align=PP_ALIGN.CENTER)

# ==================== SLIDE 4: FIND THE LETTER B ====================
s4 = add_slide(prs)
set_bg(s4, 255, 235, 220)

add_rounded_rect(s4, 0.4, 0.2, 9.2, 0.9, (220, 120, 60))
add_textbox(s4, "🔤 Find the Letter B!", 0.4, 0.22, 9.2, 0.85,
            36, bold=True, color=(255, 255, 255), align=PP_ALIGN.CENTER)
add_textbox(s4, "Circle all the letter  B  you can find!",
            0.8, 1.25, 8.4, 0.5, 18, color=(160, 80, 20), align=PP_ALIGN.CENTER)

letters = list("BCDBEFBGHBIJBKLMBNOBBPQBRSB")
per_row = 13
rows = [letters[i:i+per_row] for i in range(0, len(letters), per_row)]
for ri, row in enumerate(rows):
    for ci, letter in enumerate(row):
        x = 0.55 + ci * 0.68
        y = 1.95 + ri * 0.75
        color = (220, 50, 50) if letter == "B" else (80, 80, 80)
        add_textbox(s4, letter, x, y, 0.6, 0.6, 22, bold=(letter=="B"), color=color, align=PP_ALIGN.CENTER)

add_textbox(s4, "I found _______ letter B's!   (Hint: there are 9! 🎉)",
            1.0, 5.6, 8.0, 0.5, 15, color=(160, 80, 20), align=PP_ALIGN.CENTER)
add_textbox(s4, "🌟 Fun Learning for Little Ones  🌈  Pre-K",
            1.5, 6.9, 7.0, 0.4, 11, color=(180, 100, 40), align=PP_ALIGN.CENTER)

# ==================== SLIDE 5: ODD ONE OUT ====================
s5 = add_slide(prs)
set_bg(s5, 255, 220, 240)

add_rounded_rect(s5, 0.4, 0.2, 9.2, 0.9, (200, 80, 140))
add_textbox(s5, "🦁 Odd One Out!", 0.4, 0.22, 9.2, 0.85,
            36, bold=True, color=(255, 255, 255), align=PP_ALIGN.CENTER)
add_textbox(s5, "Circle the one that does NOT belong in each row!",
            0.8, 1.25, 8.4, 0.5, 16, color=(160, 40, 100), align=PP_ALIGN.CENTER)

rows_data = [
    ("🐶  🐱  🐭  🐻  🍕", "🍕 doesn't belong"),
    ("🍎  🍌  🍓  ✏️  🍇", "✏️ doesn't belong"),
    ("🚗  🚌  ✈️  🚂  🌸", "🌸 doesn't belong"),
    ("🔴  🔵  🟡  🟢  🐘", "🐘 doesn't belong"),
]
for i, (items, hint) in enumerate(rows_data):
    y = 1.95 + i * 1.1
    add_rounded_rect(s5, 0.6, y, 8.8, 0.88, (255, 255, 255))
    add_textbox(s5, items, 0.75, y + 0.12, 6.2, 0.65, 26)
    add_textbox(s5, hint, 7.0, y + 0.22, 2.2, 0.45, 9, italic=True, color=(180, 180, 180))

add_textbox(s5, "🌟 Fun Learning for Little Ones  🌈  Pre-K",
            1.5, 6.9, 7.0, 0.4, 11, color=(160, 60, 120), align=PP_ALIGN.CENTER)

# ==================== SLIDE 6: TRACE THE LETTERS ====================
s6 = add_slide(prs)
set_bg(s6, 240, 255, 245)

add_rounded_rect(s6, 0.4, 0.2, 9.2, 0.9, (60, 180, 120))
add_textbox(s6, "✏️ Trace the Letters!", 0.4, 0.22, 9.2, 0.85,
            36, bold=True, color=(255, 255, 255), align=PP_ALIGN.CENTER)
add_textbox(s6, "Trace over each dotted letter with your pencil or crayon!",
            0.8, 1.25, 8.4, 0.5, 16, color=(40, 140, 80), align=PP_ALIGN.CENTER)

trace_letters = ["A", "B", "C", "D"]
for i, letter in enumerate(trace_letters):
    x = 0.55 + i * 2.3
    add_rounded_rect(s6, x, 1.9, 2.1, 3.8, (255, 255, 255))
    add_textbox(s6, letter, x + 0.15, 2.0, 1.8, 2.8, 110, bold=True,
                color=(200, 200, 200), align=PP_ALIGN.CENTER)
    add_textbox(s6, f"Trace '{letter}'! ✏️", x + 0.1, 5.0, 1.9, 0.4,
                12, color=(80, 140, 100), align=PP_ALIGN.CENTER)

add_textbox(s6, "🌟 Fun Learning for Little Ones  🌈  Pre-K",
            1.5, 6.9, 7.0, 0.4, 11, color=(60, 160, 100), align=PP_ALIGN.CENTER)

# ==================== SLIDE 7: TALL OR SHORT? ====================
s7 = add_slide(prs)
set_bg(s7, 255, 250, 220)

add_rounded_rect(s7, 0.4, 0.2, 9.2, 0.9, (200, 170, 50))
add_textbox(s7, "📏 Tall or Short?", 0.4, 0.22, 9.2, 0.85,
            36, bold=True, color=(255, 255, 255), align=PP_ALIGN.CENTER)
add_textbox(s7, "Circle the TALLER one in each box!",
            0.8, 1.25, 8.4, 0.5, 18, color=(140, 110, 0), align=PP_ALIGN.CENTER)

pairs = [
    ("🦒 Giraffe", "🐶 Dog"),
    ("🌳 Tree", "🌸 Flower"),
    ("🏢 Building", "🏠 House"),
    ("👨 Dad", "👶 Baby"),
    ("⛰️ Mountain", "🏔️ Hill"),
    ("🚀 Rocket", "🚗 Car"),
]
for i, (a, b) in enumerate(pairs):
    col = i % 3
    row = i // 3
    x = 0.45 + col * 3.2
    y = 1.9 + row * 2.4
    add_rounded_rect(s7, x, y, 2.9, 2.15, (255, 255, 255))
    add_textbox(s7, a, x + 0.1, y + 0.12, 2.7, 0.65, 15, color=(60, 60, 60), align=PP_ALIGN.CENTER)
    add_textbox(s7, "VS", x + 1.05, y + 0.78, 0.8, 0.4, 12, bold=True, color=(200, 100, 0), align=PP_ALIGN.CENTER)
    add_textbox(s7, b, x + 0.1, y + 1.2, 2.7, 0.65, 15, color=(60, 60, 60), align=PP_ALIGN.CENTER)

add_textbox(s7, "🌟 Fun Learning for Little Ones  🌈  Pre-K",
            1.5, 6.9, 7.0, 0.4, 11, color=(160, 130, 0), align=PP_ALIGN.CENTER)

# ==================== SLIDE 8: CONNECT THE DOTS ====================
s8 = add_slide(prs)
set_bg(s8, 235, 220, 255)

add_rounded_rect(s8, 0.4, 0.2, 9.2, 0.9, (130, 80, 220))
add_textbox(s8, "🔢 Connect the Dots!", 0.4, 0.22, 9.2, 0.85,
            36, bold=True, color=(255, 255, 255), align=PP_ALIGN.CENTER)
add_textbox(s8, "Connect the dots from 1 to 10 to reveal the hidden picture!",
            0.8, 1.25, 8.4, 0.5, 16, color=(90, 50, 180), align=PP_ALIGN.CENTER)

dots = [
    (1, 2.8, 4.2), (2, 4.2, 2.0), (3, 5.8, 2.2),
    (4, 6.8, 3.4), (5, 6.2, 4.7), (6, 4.9, 5.2),
    (7, 3.5, 5.0), (8, 2.4, 4.2), (9, 2.2, 3.2),
    (10, 2.8, 4.2),
]
for n, x, y in dots[:9]:
    add_rounded_rect(s8, x - 0.22, y - 0.22, 0.44, 0.44, (130, 80, 220))
    add_textbox(s8, str(n), x - 0.2, y - 0.18, 0.4, 0.36, 11,
                bold=True, color=(255, 255, 255), align=PP_ALIGN.CENTER)

add_textbox(s8, "What did you draw? ____________________",
            2.2, 5.8, 5.6, 0.5, 15, color=(90, 50, 180), align=PP_ALIGN.CENTER)
add_textbox(s8, "🌟 Fun Learning for Little Ones  🌈  Pre-K",
            1.5, 6.9, 7.0, 0.4, 11, color=(110, 70, 200), align=PP_ALIGN.CENTER)

# ==================== SLIDE 9: FILL IN THE MISSING LETTER ====================
s9 = add_slide(prs)
set_bg(s9, 220, 245, 255)

add_rounded_rect(s9, 0.4, 0.2, 9.2, 0.9, (50, 150, 220))
add_textbox(s9, "🔡 Fill in the Missing Letter!", 0.4, 0.22, 9.2, 0.85,
            34, bold=True, color=(255, 255, 255), align=PP_ALIGN.CENTER)
add_textbox(s9, "Write the missing letter in each blank!",
            0.8, 1.25, 8.4, 0.5, 16, color=(30, 110, 180), align=PP_ALIGN.CENTER)

sequences = [
    "A   B   ___   D   E",
    "F   ___   H   I   J",
    "K   L   M   ___   O",
    "P   Q   ___   S   T",
    "U   V   W   ___   Y",
]
for i, seq in enumerate(sequences):
    y = 1.95 + i * 0.95
    add_rounded_rect(s9, 1.0, y, 8.0, 0.75, (255, 255, 255))
    add_textbox(s9, seq, 1.1, y + 0.12, 7.8, 0.52, 22, bold=True,
                color=(50, 50, 150), align=PP_ALIGN.CENTER)

add_textbox(s9, "🌟 Fun Learning for Little Ones  🌈  Pre-K",
            1.5, 6.9, 7.0, 0.4, 11, color=(40, 120, 190), align=PP_ALIGN.CENTER)

# ==================== SLIDE 10: WHAT SOUND? ====================
s10 = add_slide(prs)
set_bg(s10, 255, 240, 220)

add_rounded_rect(s10, 0.4, 0.2, 9.2, 0.9, (220, 140, 40))
add_textbox(s10, "🔊 What Sound Does It Start With?", 0.4, 0.22, 9.2, 0.85,
            30, bold=True, color=(255, 255, 255), align=PP_ALIGN.CENTER)
add_textbox(s10, "Draw a line from the picture to the letter it starts with!",
            0.8, 1.25, 8.4, 0.5, 16, color=(160, 90, 20), align=PP_ALIGN.CENTER)

word_pairs = [
    ("🍎  Apple",  "S"),
    ("⚽  Ball",   "A"),
    ("🐱  Cat",    "B"),
    ("🐶  Dog",    "C"),
    ("☀️  Sun",    "D"),
]
for i, (pic, letter) in enumerate(word_pairs):
    y = 1.9 + i * 0.92
    add_rounded_rect(s10, 0.8, y, 2.6, 0.72, (255, 255, 255))
    add_textbox(s10, pic, 0.9, y + 0.12, 2.4, 0.5, 16, color=(60, 60, 60))
    add_rounded_rect(s10, 6.8, y, 1.3, 0.72, (255, 255, 255))
    add_textbox(s10, letter, 6.85, y + 0.08, 1.2, 0.55, 22, bold=True,
                color=(180, 80, 0), align=PP_ALIGN.CENTER)

add_textbox(s10, "🌟 Fun Learning for Little Ones  🌈  Pre-K",
            1.5, 6.9, 7.0, 0.4, 11, color=(180, 110, 30), align=PP_ALIGN.CENTER)

# ==================== SLIDE 11: CERTIFICATE ====================
s11 = add_slide(prs)
set_bg(s11, 255, 252, 220)

add_rounded_rect(s11, 0.3, 0.2, 9.4, 7.0, (255, 240, 180))
add_rounded_rect(s11, 0.55, 0.42, 8.9, 6.6, (255, 255, 255))

add_textbox(s11, "🌟  Super Brain Booster!  🌟", 0.7, 0.55, 8.6, 0.9,
            36, bold=True, color=(200, 150, 0), align=PP_ALIGN.CENTER)
add_textbox(s11, "✨ ✨ ✨ ✨ ✨ ✨ ✨ ✨ ✨ ✨", 1.0, 1.48, 8.0, 0.4,
            14, color=(255, 200, 0), align=PP_ALIGN.CENTER)
add_textbox(s11, "This certificate is proudly awarded to:", 1.5, 1.95, 7.0, 0.45,
            16, color=(100, 100, 100), align=PP_ALIGN.CENTER)
add_textbox(s11, "____________________________", 2.0, 2.45, 6.0, 0.55,
            20, color=(80, 80, 80), align=PP_ALIGN.CENTER)
add_textbox(s11, "for completing all Brain Booster activities! 🎉🎊", 1.2, 3.1, 7.6, 0.5,
            15, color=(100, 100, 100), align=PP_ALIGN.CENTER)

add_textbox(s11, "My Scores 📊", 3.5, 3.7, 3.0, 0.45,
            18, bold=True, color=(150, 100, 0), align=PP_ALIGN.CENTER)

score_items = ["Counting 🍎", "Shapes 🔷", "Letters 🔤", "Patterns 🔢", "Tracing ✏️", "Sounds 🔊"]
for i, item in enumerate(score_items):
    col = i % 3
    row = i // 3
    x = 0.75 + col * 3.0
    y = 4.25 + row * 0.88
    add_rounded_rect(s11, x, y, 2.7, 0.72, (255, 245, 200))
    add_textbox(s11, item, x + 0.1, y + 0.06, 2.5, 0.3, 12, color=(100, 80, 0), align=PP_ALIGN.CENTER)
    add_textbox(s11, "⭐ ⭐ ⭐ ⭐ ⭐", x + 0.1, y + 0.38, 2.5, 0.28, 10,
                color=(200, 200, 200), align=PP_ALIGN.CENTER)

add_textbox(s11, "🎉  Amazing Work! You're a Brain Booster Star!  🎉",
            0.9, 6.35, 8.2, 0.55, 15, bold=True, color=(180, 130, 0), align=PP_ALIGN.CENTER)

# ==================== SAVE ====================
out_path = r"C:\Users\hp\OneDrive\Desktop\Kids Playbook\Brain_Boosters_PreK.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
